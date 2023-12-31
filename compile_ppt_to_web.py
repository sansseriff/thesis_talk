import collections
import collections.abc
import pptx
import win32com.client as client
import os
import shutil
import vss
from pyuac import main_requires_admin
import time
import subprocess
import argparse
import json



def prev_slide(number):
    if number > 0:
        return number - 1
    else:
        return 0


def next_slide(number, ln):
    if number < ln - 1:
        return number + 1
    else:
        return ln - 1
    

# old

# <script>
#   import {{ onMount }} from 'svelte';
#   import {{ writable }} from 'svelte/store';
#   import {{ goto }} from '$app/navigation';
#   // <img src="/slides_png/slide_{i}.png" alt="slide_{i}" width="1000" height="500">

#   const width = writable(0);
#   const height = writable(0);
#   onMount(() => {{
#     width.set(window.innerWidth);
#     height.set(window.innerHeight);
#     document.addEventListener('keydown', async function(event) {{
#       //Check if the pressed key is the left arrow key
#       if (event.key === 'ArrowLeft' || event.key === 'PageUp') {{
#         // Navigate to the desired URL when the left arrow key is pressed
#         //window.location.href = '/slide_{prev_slide(i)}'; // Replace with your desired URL
#         await goto('/slide_{prev_slide(i)}');
#       }}
#       //Check if the pressed key is the right arrow key
#       else if (event.key === 'ArrowRight' || event.key === 'PageDown') {{
#         // Navigate to the desired URL when the right arrow key is pressed
#         //window.location.href = '/slide_{next_slide(i, ln)}'; // Replace with your desired URL
#         await goto('/slide_{next_slide(i, ln)}');
#       }}
#     }});
#   }});
# </script>

def snake_to_camel(snake_str):
    components = snake_str.split('_')
    return ''.join(x.title() for x in components)


def create_overlay_interactive_svelte_route(i, ln, routes_directory, component_file_name):
    # I should add a checker here that looks in the route directory
    # if its there and does not overwrite it if theres some keyfile 
    # indicating dynamic content

    folder_name = f"slide_{i}"
    folder_path = os.path.join(routes_directory, folder_name)
    os.makedirs(folder_path, exist_ok=True)


    component_name = snake_to_camel(component_file_name)

    # Create a text file in the folder
    svelte_path = os.path.join(folder_path, "+page.svelte")
    with open(svelte_path, "w") as f:
        f.write(
            f"""
<script lang='ts'>
  import {{ onMount, onDestroy }} from 'svelte';
  import {{ writable }} from 'svelte/store';
  import {{ goto }} from '$app/navigation';
  import {component_name} from '$lib/{component_file_name}.svelte';

  const width = writable(0);
  const height = writable(0);
  onMount(() => {{
    width.set(window.innerWidth);
    height.set(window.innerHeight);
  }});
</script>

<div id="floating-number">{i}</div>
<{component_name} prev_slide='/slide_{prev_slide(i)}' next_slide='/slide_{next_slide(i, ln)}'/>
<img src="/slides_png/slide_{i}.png" alt="slide_{i}" width="{{$width}}" style="top: 0px; position: absolute; z-index: -1;">
              
"""
        )

    js_path = os.path.join(folder_path, "+page.js")
    with open(js_path, "w") as f:
        f.write(
            """
// since there's no dynamic data here, we can prerender
// it so that it gets served as a static asset in production
export const prerender = true;
"""
        )
        # end



def create_overlay_video_svelte_route(i, ln, routes_directory, video_name):
    # I should add a checker here that looks in the route directory
    # if its there and does not overwrite it if theres some keyfile 
    # indicating dynamic content

    with open("video_sizing.json", "r") as f:
        video_sizing = json.loads(f.read())

    this_video_size = video_sizing.get(video_name, None)
    if this_video_size is None:
        this_video_size = video_sizing["default"]


        video_sizing[video_name] = this_video_size

    folder_name = f"slide_{i}"
    folder_path = os.path.join(routes_directory, folder_name)
    os.makedirs(folder_path, exist_ok=True)


    with open("video_sizing.json", "w") as f:
        f.write(json.dumps(video_sizing, indent=4))

    # Create a text file in the folder
    svelte_path = os.path.join(folder_path, "+page.svelte")
    with open(svelte_path, "w") as f:
        f.write(
            f"""
<script lang='ts'>
  // {video_name}
  import {{ onMount, onDestroy }} from 'svelte';
  import {{ writable }} from 'svelte/store';
  import {{ goto }} from '$app/navigation';

  const width = writable(0);
  const height = writable(0);
  const video_width = writable(0);
  const video_top = writable(0);
  const video_left = writable(0);

  let keydownHandler: (event: KeyboardEvent) => void;

  onMount(() => {{
    width.set(window.innerWidth);
    height.set(window.innerHeight);

    video_width.set(Math.floor({this_video_size["width"]} * window.innerWidth))
    video_top.set(Math.floor({this_video_size["top"]} * window.innerWidth))
    video_left.set(Math.floor({this_video_size["left"]} * window.innerWidth))


    const video = document.getElementById('myVideo');
    video.play();

    keydownHandler = async function(event) {{
      // Define an async function
      const navigate = async (url: string, imgSrc: string) => {{
        const img = new Image();
        img.src = imgSrc;
        img.onload = async () => {{
          await goto(url);
        }};
      }}
      //Check if the pressed key is the left arrow key
      if (event.key === 'ArrowLeft' || event.key === 'PageUp') {{
        // Navigate to the desired URL when the left arrow key is pressed
        navigate('/slide_{prev_slide(i)}', '/slides_png/slide_{prev_slide(i)}.png');
      }}
      //Check if the pressed key is the right arrow key
      else if (event.key === 'ArrowRight' || event.key === 'PageDown') {{
        // Navigate to the desired URL when the right arrow key is pressed
        navigate('/slide_{next_slide(i, ln)}', '/slides_png/slide_{next_slide(i, ln)}.png');
      }}
    }};

    document.addEventListener('keydown', keydownHandler);
  }});

  onDestroy(() => {{
    if (typeof window !== 'undefined') {{
      document.removeEventListener('keydown', keydownHandler);
    }}
  }});
</script>

<div id="floating-number">{i}</div>
<video id="myVideo" style="position:absolute; top:{{$video_top}}px; left:{{$video_left}}px; width:{{$video_width}}px;" src="/slides_videos/{video_name}.mp4" muted playsinline></video>
<img src="/slides_png/slide_{i}.png" alt="slide_{i}" width="{{$width}}">

              
"""
        )

    js_path = os.path.join(folder_path, "+page.js")
    with open(js_path, "w") as f:
        f.write(
            """
// since there's no dynamic data here, we can prerender
// it so that it gets served as a static asset in production
export const prerender = true;
"""
        )

  
        # end


####################################

def create_svelte_route(i, ln, routes_directory):
    # I should add a checker here that looks in the route directory
    # if its there and does not overwrite it if theres some keyfile 
    # indicating dynamic content

    folder_name = f"slide_{i}"
    folder_path = os.path.join(routes_directory, folder_name)
    os.makedirs(folder_path, exist_ok=True)

    # Create a text file in the folder
    svelte_path = os.path.join(folder_path, "+page.svelte")
    with open(svelte_path, "w") as f:
        f.write(
            f"""
<script lang='ts'>
  import {{ onMount, onDestroy }} from 'svelte';
  import {{ writable }} from 'svelte/store';
  import {{ goto }} from '$app/navigation';

  const width = writable(0);
  const height = writable(0);

  let keydownHandler: (event: KeyboardEvent) => void;

  onMount(() => {{
    width.set(window.innerWidth);
    height.set(window.innerHeight);

    keydownHandler = async function(event) {{
      // Define an async function
      const navigate = async (url: string, imgSrc: string) => {{
        const img = new Image();
        img.src = imgSrc;
        img.onload = async () => {{
          await goto(url);
        }};
      }}
      //Check if the pressed key is the left arrow key
      if (event.key === 'ArrowLeft' || event.key === 'PageUp') {{
        // Navigate to the desired URL when the left arrow key is pressed
        navigate('/slide_{prev_slide(i)}', '/slides_png/slide_{prev_slide(i)}.png');
      }}
      //Check if the pressed key is the right arrow key
      else if (event.key === 'ArrowRight' || event.key === 'PageDown') {{
        // Navigate to the desired URL when the right arrow key is pressed
        navigate('/slide_{next_slide(i, ln)}', '/slides_png/slide_{next_slide(i, ln)}.png');
      }}
    }};

    document.addEventListener('keydown', keydownHandler);
  }});

  onDestroy(() => {{
    if (typeof window !== 'undefined') {{
      document.removeEventListener('keydown', keydownHandler);
    }}
  }});
</script>

<div id="floating-number">{i}</div>
<img src="/slides_png/slide_{i}.png" alt="slide_{i}" width="{{$width}}">
              
"""
        )

    js_path = os.path.join(folder_path, "+page.js")
    with open(js_path, "w") as f:
        f.write(
            """
// since there's no dynamic data here, we can prerender
// it so that it gets served as a static asset in production
export const prerender = true;
"""
        )
        # end



def create_backup_supporting_svelte_route(i, ln, routes_directory, b):
    # I should add a checker here that looks in the route directory
    # if its there and does not overwrite it if theres some keyfile 
    # indicating dynamic content
    if b == 0:
        folder_name = f"slide_{i}"
    else:
        folder_name = f"slide_{i}_backup_{b}"

    print("folder_name: ", folder_name)
    folder_path = os.path.join(routes_directory, folder_name)
    os.makedirs(folder_path, exist_ok=True)

    # Create a text file in the folder
    svelte_path = os.path.join(folder_path, "+page.svelte")
    with open(svelte_path, "w") as f:
        f.write(
            f"""
<script lang='ts'>
  import {{ onMount, onDestroy }} from 'svelte';
  import {{ writable }} from 'svelte/store';
  import {{ goto }} from '$app/navigation';

  const width = writable(0);
  const height = writable(0);

  let keydownHandler: (event: KeyboardEvent) => void;

  onMount(() => {{
    width.set(window.innerWidth);
    height.set(window.innerHeight);

    keydownHandler = async function(event) {{
      // Define an async function
      const navigate = async (url: string, imgSrc: string) => {{
        const img = new Image();
        img.src = imgSrc;
        img.onload = async () => {{
          await goto(url);
        }};
      }}
      //Check if the pressed key is the left arrow key
      if (event.key === 'ArrowLeft' || event.key === 'PageUp') {{
        // Navigate to the desired URL when the left arrow key is pressed
        navigate('/slide_{prev_slide(i)}', '/slides_png/slide_{prev_slide(i)}.png');
      }}
      //Check if the pressed key is the right arrow key
      else if (event.key === 'ArrowRight' || event.key === 'PageDown') {{
        // Navigate to the desired URL when the right arrow key is pressed
        navigate('/slide_{next_slide(i, ln)}', '/slides_png/slide_{next_slide(i, ln)}.png');
      }}
""")
    # handle the backup
    with open(svelte_path, "a") as f:
        f.write(
            f"""
      else if (event.key === '.') {{
        // Navigate to the backup slide
        navigate('/slide_{i}_backup_{b+1}', '/slides_png/slide_{i}_backup_{b+1}.png');
      }}
    }};
""")
        
    if b == 0:
        q = i
    else:
        q = f"{i}_backup_{b}"

    with open(svelte_path, "a") as f:
          f.write(f"""


    document.addEventListener('keydown', keydownHandler);
  }});

  onDestroy(() => {{
    if (typeof window !== 'undefined') {{
      document.removeEventListener('keydown', keydownHandler);
    }}
  }});
</script>

<div id="floating-number">{i}</div>
<img src="/slides_png/slide_{q}.png" alt="slide_{q}" width="{{$width}}">
              
"""
        )

    js_path = os.path.join(folder_path, "+page.js")
    with open(js_path, "w") as f:
        f.write(
            """
// since there's no dynamic data here, we can prerender
// it so that it gets served as a static asset in production
export const prerender = true;
"""
        )
        # end




file_export = "./powerpoint/defense_export.pptx"
export_path = ".\\web\\defense\\static\\slides_png\\"
routes_directory = "./web/defense/src/routes/"
full_export_path = os.path.abspath(export_path)

full_file_path = os.path.abspath(file_export)

def run_copy():
    subprocess.run(["python", "copier.py"])

def finish_up(run_copy_flag=True):
    # if run_copy_flag:
    run_copy()
        
    


    # works
    # ppt=Presentation("/Users/Andrew/defense.pptx")
    





    # Define the base directory where you want to create the folders
    

    # Define the number of folders you want to create

    ppt = pptx.Presentation(file_export)
    notes = []
    # https://stackoverflow.com/questions/61815883/python-pptx-export-img-png-jpeg

    Application = client.Dispatch("PowerPoint.Application")
    Presentation = Application.Presentations.Open(full_file_path)

    j = 0
    backup_counter = 0
    # Loop through the range of folder numbers and create each folder
    for i, (note_slide, img_slide) in enumerate(zip(ppt.slides, Presentation.Slides)):
        textNote = note_slide.notes_slide.notes_text_frame.text
        notes.append((i, textNote))

        if "[interactive][overlay]" in textNote:
            img_slide.Export(f"{full_export_path}\\slide_{j}.png", "PNG")
            split_string = textNote.split("[interactive][overlay]", 1)
            component_file_name = split_string[1].replace("[", "").replace("]", "")
            component_file_name = component_file_name.split("#", 1)[0]
            print("found interactive slide for filename: ", component_file_name)
            create_overlay_interactive_svelte_route(j, len(ppt.slides), routes_directory, component_file_name)
            j += 1
            continue

        if "[dynamic][overlay]" in textNote:
            img_slide.Export(f"{full_export_path}\\slide_{j}.png", "PNG")

            split_string = textNote.split("[dynamic][overlay]", 1)
            video_name = split_string[1].replace("[", "").replace("]", "")
            video_name = video_name.split("#", 1)[0]
            print("found dynamic slide for: ", video_name)
            create_overlay_video_svelte_route(j, len(ppt.slides), routes_directory, video_name)
            j += 1
            continue

        if ("[backup_root]" in textNote) or ("[backup]" in textNote):
            # handle slide that supports backup
            print("found backup slide, i=", i, "j=", j)
            if backup_counter == 0:
              img_slide.Export(f"{full_export_path}\\slide_{j}.png", "PNG")
            else:
              img_slide.Export(f"{full_export_path}\\slide_{j}_backup_{backup_counter}.png", "PNG")
            create_backup_supporting_svelte_route(j, len(ppt.slides), routes_directory, backup_counter)
            # if backup_counter == 0:
            #     j += 1
            backup_counter += 1
            continue
            
        # don't overwrite the slide that supports backup
        if backup_counter > 0:
            j += 1
        backup_counter = 0
          
        # handle normal slide
        # if "[dynamic]" not in textNote:
        img_slide.Export(f"{full_export_path}\\slide_{j}.png", "PNG")
        
        create_svelte_route(j, len(ppt.slides), routes_directory)
        j += 1

        
        

              

    with open(f"{full_export_path}\\notes.json", "w") as f:
        f.write(str(notes))

    

    

    # Application.Quit() # this closes all instances of powerpoint.
    Presentation.Close() # this just closes the auto-opened instance
    Presentation = None
    Application = None


def fix_svelte_routes():
    with open(f"{full_export_path}\\notes.json", "r") as f:
        notes = eval(f.read())

    # print(notes)

    for i, note in notes:
        if "dyno" not in note:
            create_svelte_route(i, len(notes), routes_directory)
        else:
            print("found dynamic slide")


if __name__ == "__main__":

    parser = argparse.ArgumentParser()
    parser.add_argument("--no-copy", action="store_true", help="Deactivate the run_copy function")
    args = parser.parse_args()

    # print(not args.no_copy)

    # finish_up(not args.no_copy)
    if not args.no_copy:
        finish_up()
    else:
        fix_svelte_routes()

